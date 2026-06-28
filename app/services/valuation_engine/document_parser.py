from __future__ import annotations

from pathlib import Path
from typing import Any


SHORT_TEXT_WARNING = "上传的资料可能是扫描件或图片型材料，当前版本可能无法完整识别。"
SCAN_WARNING = "该 PDF 可能为扫描件，当前版本可能无法完整识别。后续可接入 OCR。"
PPT_LEGACY_WARNING = "当前建议使用 PPTX 格式，旧版 PPT 可能无法完整解析。"
DOC_LEGACY_WARNING = "当前建议使用 DOCX 格式，旧版 DOC 可能无法完整解析。"


def parse_uploaded_document(file_path: str | Path) -> dict[str, Any]:
    path = Path(file_path).expanduser()
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return parse_pdf(path)
        if suffix == ".pptx":
            return parse_pptx(path)
        if suffix == ".ppt":
            return parse_ppt_legacy_warning(path)
        if suffix == ".docx":
            return parse_docx(path)
        if suffix == ".doc":
            return parse_doc_legacy_warning(path)
        return base_result(path, "unknown", "unsupported", [f"暂不支持 {suffix or '未知'} 文件类型。"], "failed")
    except Exception as exc:  # pragma: no cover - last line of defense for Streamlit
        return base_result(path, suffix.lstrip(".") or "unknown", "failed", [f"文件解析失败：{exc}"], "failed")


def parse_project_document(file_path: str | Path) -> dict[str, Any]:
    return parse_uploaded_document(file_path)


def parse_pdf(file_path: str | Path) -> dict[str, Any]:
    path = Path(file_path).expanduser()
    warnings: list[str] = []
    pages: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []

    try:
        import fitz  # type: ignore
    except ImportError:
        fitz = None
        warnings.append("未安装 PyMuPDF，PDF 正文提取能力受限。")

    if fitz is not None:
        try:
            with fitz.open(path) as doc:
                for index, page in enumerate(doc, start=1):
                    pages.append({"page_number": index, "text": (page.get_text("text") or "").strip()})
        except Exception as exc:
            warnings.append(f"PyMuPDF 读取失败：{exc}")

    try:
        import pdfplumber  # type: ignore
    except ImportError:
        pdfplumber = None
        warnings.append("未安装 pdfplumber，PDF 表格提取能力受限。")

    if pdfplumber is not None:
        try:
            with pdfplumber.open(path) as pdf:
                if not pages:
                    pages = [
                        {"page_number": index, "text": (page.extract_text() or "").strip()}
                        for index, page in enumerate(pdf.pages, start=1)
                    ]
                for index, page in enumerate(pdf.pages, start=1):
                    for table_index, table in enumerate(page.extract_tables() or [], start=1):
                        tables.append({"page_number": index, "table_index": table_index, "rows": table})
        except Exception as exc:
            warnings.append(f"pdfplumber 读取失败：{exc}")

    raw_text = "\n\n".join(page["text"] for page in pages if page.get("text")).strip()
    if len(raw_text) < 120:
        warnings.extend([SCAN_WARNING, SHORT_TEXT_WARNING])

    return {
        **base_result(path, "pdf", "PyMuPDF + pdfplumber", warnings, quality_from_text(raw_text, pages, tables)),
        "pages": pages,
        "raw_text": raw_text,
        "tables": tables,
    }


def parse_pptx(file_path: str | Path) -> dict[str, Any]:
    path = Path(file_path).expanduser()
    warnings: list[str] = []
    slides: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return base_result(path, "pptx", "python-pptx", ["未安装 python-pptx，无法解析 PPTX。"], "failed")

    try:
        presentation = Presentation(path)
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = ""
            body_parts: list[str] = []
            slide_tables: list[dict[str, Any]] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = clean_text(getattr(shape, "text", ""))
                    if text:
                        if not title:
                            title = text.splitlines()[0]
                        body_parts.append(text)
                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        rows.append([clean_text(cell.text) for cell in row.cells])
                    table = {"slide_number": slide_number, "rows": rows}
                    slide_tables.append(table)
                    tables.append(table)
            notes = ""
            try:
                notes = clean_text(slide.notes_slide.notes_text_frame.text)
            except Exception:
                notes = ""
            slides.append(
                {
                    "slide_number": slide_number,
                    "title": title,
                    "text": "\n".join(body_parts).strip(),
                    "notes": notes,
                    "tables": slide_tables,
                }
            )
    except Exception as exc:
        return base_result(path, "pptx", "python-pptx", [f"PPTX 读取失败：{exc}"], "failed")

    raw_text_parts = []
    for slide in slides:
        table_text = "\n".join(
            " | ".join(cell for cell in row if cell)
            for table in slide.get("tables", [])
            for row in table.get("rows", [])
        )
        raw_text_parts.append(
            "\n".join(part for part in [slide.get("title", ""), slide.get("text", ""), slide.get("notes", ""), table_text] if part)
        )
    raw_text = "\n\n".join(raw_text_parts).strip()
    if len(raw_text) < 80:
        warnings.append(SHORT_TEXT_WARNING)

    return {
        **base_result(path, "pptx", "python-pptx", warnings, quality_from_text(raw_text, slides, tables)),
        "slides": slides,
        "tables": tables,
        "raw_text": raw_text,
    }


def parse_ppt_legacy_warning(file_path: str | Path) -> dict[str, Any]:
    path = Path(file_path).expanduser()
    return base_result(path, "ppt", "legacy_warning", [PPT_LEGACY_WARNING], "failed")


def parse_docx(file_path: str | Path) -> dict[str, Any]:
    path = Path(file_path).expanduser()
    warnings: list[str] = []
    paragraphs: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []

    try:
        from docx import Document  # type: ignore
    except ImportError:
        return base_result(path, "docx", "python-docx", ["未安装 python-docx，无法解析 DOCX。"], "failed")

    try:
        document = Document(path)
        for index, paragraph in enumerate(document.paragraphs, start=1):
            text = clean_text(paragraph.text)
            if text:
                paragraphs.append({"paragraph_number": index, "style": paragraph.style.name if paragraph.style else "", "text": text})
        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                rows.append([clean_text(cell.text) for cell in row.cells])
            tables.append({"table_index": table_index, "rows": rows})
    except Exception as exc:
        return base_result(path, "docx", "python-docx", [f"DOCX 读取失败：{exc}"], "failed")

    raw_text = "\n".join(paragraph["text"] for paragraph in paragraphs)
    table_text = "\n".join(" | ".join(cell for cell in row if cell) for table in tables for row in table.get("rows", []))
    raw_text = "\n\n".join(part for part in [raw_text, table_text] if part).strip()
    if len(raw_text) < 80:
        warnings.append(SHORT_TEXT_WARNING)

    return {
        **base_result(path, "docx", "python-docx", warnings, quality_from_text(raw_text, paragraphs, tables)),
        "paragraphs": paragraphs,
        "tables": tables,
        "raw_text": raw_text,
    }


def parse_doc_legacy_warning(file_path: str | Path) -> dict[str, Any]:
    path = Path(file_path).expanduser()
    return base_result(path, "doc", "legacy_warning", [DOC_LEGACY_WARNING], "failed")


def base_result(
    path: Path,
    file_type: str,
    parser: str,
    warnings: list[str] | None = None,
    extraction_quality: str = "failed",
) -> dict[str, Any]:
    return {
        "file_name": path.name,
        "file_path": str(path),
        "file_type": file_type,
        "parser": parser,
        "raw_text": "",
        "pages": [],
        "slides": [],
        "paragraphs": [],
        "tables": [],
        "warnings": dedupe(warnings or []),
        "extraction_quality": extraction_quality,
    }


def quality_from_text(raw_text: str, sections: list[Any], tables: list[Any]) -> str:
    if not raw_text.strip():
        return "failed"
    if len(raw_text) >= 1200 or (len(raw_text) >= 400 and (len(sections) >= 3 or tables)):
        return "high"
    if len(raw_text) >= 120 or (len(raw_text) >= 40 and (sections or tables)):
        return "medium"
    return "low"


def clean_text(value: str) -> str:
    return "\n".join(line.strip() for line in str(value or "").splitlines() if line.strip())


def dedupe(items: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        if item and item not in seen:
            seen.add(item)
            result.append(item)
    return result
